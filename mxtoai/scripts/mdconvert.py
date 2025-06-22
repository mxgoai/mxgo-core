# This is copied from Magentic-one's great repo: https://github.com/microsoft/autogen/blob/v0.4.4/python/packages/autogen-magentic-one/src/autogen_magentic_one/markdown_browser/mdconvert.py
# Thanks to Microsoft researchers for open-sourcing this!
# type: ignore
import base64
import contextlib
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
from typing import Any, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx

# File-format detection
import pydub
import requests
import speech_recognition as sr
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline and not re.search(r"^\n", text):
            return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool = False, **kwargs):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return f"{prefix}{text}{suffix}"
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return f"{prefix}{text}{suffix}"

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return f"<{href}>"
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "{}"'.format(title.replace('"', r"\"")) if title else ""
        return f"{prefix}[{text}]({href}{title_part}){suffix}" if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool = False, **kwargs) -> str:
        """Same as usual converter, but removes data URIs"""
        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "{}"'.format(title.replace('"', r"\"")) if title else ""
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return f"![{alt}]({src}{title_part})"

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        """Convert content from memory. Must be implemented by subclasses."""
        raise NotImplementedError


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        content_type, _ = mimetypes.guess_type("__placeholder" + kwargs.get("file_extension", ""))

        if content_type is None or not content_type.startswith("text/"):
            return None

        if b"\0" in content[:1024]:
            return None

        try:
            text_content = content.decode("utf-8")
            return DocumentConverterResult(
                title=None,
                text_content=text_content,
            )
        except UnicodeDecodeError:
            return None


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        try:
            html_content = content.decode("utf-8")
            return self._convert(html_content)
        except UnicodeDecodeError:
            return None

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""
        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string, text_content=webpage_text
        )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        try:
            html_content = content.decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
        except UnicodeDecodeError:
            return None

        for script in soup(["script", "style"]):
            script.extract()

        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        try:
            html_content = content.decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
        except UnicodeDecodeError:
            return None

        assert soup.title is not None
        assert soup.title.string is not None
        metadata: dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        transcript_text = ""
        parsed_url = urlparse(url)  # type: ignore
        params = parse_qs(parsed_url.query)  # type: ignore
        if "v" in params:
            assert isinstance(params["v"][0], str)
            video_id = str(params["v"][0])
            try:
                # Must be a single transcript.
                transcript = YouTubeTranscriptApi.get_transcript(video_id)  # type: ignore
                # Alternative formatting:
                transcript_text = SRTFormatter().format_transcript(transcript)
            except Exception:
                pass
        if transcript_text:
            webpage_text += f"\n### Transcript\n{transcript_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(self, metadata: dict[str, str], keys: list[str], default: Union[str, None] = None) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(self, json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                ret = self._findKey(json[k], key)
                if ret is not None:
                    return ret
        return None


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            text_content = pdfminer.high_level.extract_text(temp_file_path)
            os.unlink(temp_file_path)

            return DocumentConverterResult(
                title=None,
                text_content=text_content,
            )
        except Exception:
            return None


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        try:
            from io import BytesIO

            result = mammoth.convert_to_html(BytesIO(content))
            html_content = result.value
            return self._convert(html_content)
        except Exception:
            return None


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        try:
            from io import BytesIO

            sheets = pd.read_excel(BytesIO(content), sheet_name=None)
            md_content = ""
            for s in sheets:
                md_content += f"## {s}\n"
                html_content = sheets[s].to_html(index=False)
                md_content += self._convert(html_content).text_content.strip() + "\n\n"

            return DocumentConverterResult(
                title=None,
                text_content=md_content.strip(),
            )
        except Exception:
            return None


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        try:
            from io import BytesIO

            md_content = ""

            presentation = pptx.Presentation(BytesIO(content))
            slide_num = 0
            for slide in presentation.slides:
                slide_num += 1

                md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

                title = slide.shapes.title
                for shape in slide.shapes:
                    if self._is_picture(shape):
                        alt_text = ""
                        with contextlib.suppress(Exception):
                            alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

                    if self._is_table(shape):
                        html_table = "<html><body><table>"
                        first_row = True
                        for row in shape.table.rows:
                            html_table += "<tr>"
                            for cell in row.cells:
                                if first_row:
                                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                                else:
                                    html_table += "<td>" + html.escape(cell.text) + "</td>"
                            html_table += "</tr>"
                            first_row = False
                        html_table += "</table></body></html>"
                        md_content += "\n" + self._convert(html_table).text_content.strip() + "\n"

                    elif shape.has_text_frame:
                        if shape == title:
                            md_content += "# " + shape.text.lstrip() + "\n"
                        else:
                            md_content += shape.text + "\n"

                md_content = md_content.strip()

                if slide.has_notes_slide:
                    md_content += "\n\n### Notes:\n"
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame is not None:
                        md_content += notes_frame.text
                    md_content = md_content.strip()

            return DocumentConverterResult(
                title=None,
                text_content=md_content.strip(),
            )
        except Exception:
            return None

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        return bool(shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"))

    def _is_table(self, shape):
        return shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        try:
            result = subprocess.run([exiftool, "-json", local_path], capture_output=True, text=True, check=False).stdout
            return json.loads(result)[0]
        except Exception:
            return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata and speech transcription.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        try:
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            try:
                transcript = self._transcribe_audio(temp_file_path)
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            os.unlink(temp_file_path)
        except Exception:
            return None

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 and M4A files to markdown via extraction of metadata and speech transcription.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = ""

        handle, temp_path = tempfile.mkstemp(suffix=".wav")
        os.close(handle)
        try:
            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_input:
                temp_input.write(content)
                temp_input_path = temp_input.name

            if extension.lower() == ".mp3":
                sound = pydub.AudioSegment.from_mp3(temp_input_path)
            else:
                sound = pydub.AudioSegment.from_file(temp_input_path, format="m4a")
            sound.export(temp_path, format="wav")

            try:
                transcript = super()._transcribe_audio(temp_path).strip()
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            os.unlink(temp_input_path)
        except Exception:
            return None
        finally:
            with contextlib.suppress(OSError):
                os.unlink(temp_path)

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ZipConverter(DocumentConverter):
    """
    Extracts ZIP files and returns a listing of extracted files.
    """

    def __init__(self, extract_dir: str = "downloads"):
        self.extract_dir = extract_dir
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        try:
            from io import BytesIO

            extracted_files = []

            with zipfile.ZipFile(BytesIO(content), "r") as zip_ref:
                zip_ref.extractall(self.extract_dir)
                for file_path in zip_ref.namelist():
                    if not file_path.endswith("/"):
                        extracted_files.append(self.extract_dir + "/" + file_path)

            extracted_files.sort()

            md_content = "Downloaded the following files:\n"
            for file in extracted_files:
                md_content += f"* {file}\n"

            return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())
        except Exception:
            return None


class ImageConverter(MediaConverter):
    """
    Converts images to Markdown, with optional vision-language model descriptions.
    """

    def convert_content(self, content: bytes, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"]:
            return None

        md_content = ""
        filename = kwargs.get("filename", "image" + extension)

        description = None
        client = kwargs.get("mlm_client")
        model = kwargs.get("mlm_model")
        if client and model:
            try:
                with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_file:
                    temp_file.write(content)
                    temp_file_path = temp_file.name

                description = self._get_mlm_description(temp_file_path, extension, client, model)
                os.unlink(temp_file_path)
            except Exception:
                pass

        if description:
            md_content += f"## Image Description\n{description}\n"

        md_content += f"\n![{filename}]({filename})\n"

        return DocumentConverterResult(
            title=filename,
            text_content=md_content.strip(),
        )

    def _get_mlm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class FileConversionException(Exception):
    pass


class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """
    Converts various file types to Markdown from memory content.
    """

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
    ):
        self._requests_session = requests_session if requests_session is not None else requests.Session()
        self._mlm_client = mlm_client
        self._mlm_model = mlm_model
        self._page_converters: list[DocumentConverter] = []

        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(PdfConverter())

    def convert_content(self, content: bytes, **kwargs: Any) -> DocumentConverterResult:
        """Convert content from memory bytes."""
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        filename = kwargs.get("filename", "")
        if filename:
            base, file_ext = os.path.splitext(filename)
            self._append_ext(extensions, file_ext)

        return self._convert_content(content, extensions, **kwargs)

    def _convert_content(self, content: bytes, extensions: list[Union[str, None]], **kwargs) -> DocumentConverterResult:
        error_trace = ""
        for ext in [*extensions, None]:
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                if "mlm_model" not in _kwargs and self._mlm_model is not None:
                    _kwargs["mlm_model"] = self._mlm_model

                res = None
                try:
                    res = converter.convert_content(content, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                    return res

        if len(error_trace) > 0:
            filename = kwargs.get("filename", "content")
            msg = f"Could not convert '{filename}' to Markdown. File type was recognized as {extensions}. While converting the content, the following error was encountered:\n\n{error_trace}"
            raise FileConversionException(msg)

        filename = kwargs.get("filename", "content")
        msg = f"Could not convert '{filename}' to Markdown. The formats {extensions} are not supported."
        raise UnsupportedFormatException(msg)

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        extensions.append(ext)

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
